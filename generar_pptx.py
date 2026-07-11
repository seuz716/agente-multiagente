#!/usr/bin/env python3
"""Generador PPTX con tema profesional"""
from pptx import Presentation
from pptx.util import Pt, Inches, Cm
from pptx.dml.color import RGBColor

def crear_pptx(titulo="Sistema Agente-Multiagente"):
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    
    # Slide 1: Título
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = titulo
    slide.placeholders[1].text = "Análisis y Mejoras\nJulio 2026"
    
    # Slide 2: Introducción
    slide = prs.slides.add_slide(blank_layout)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title.text = "Introducción"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(5))
    tf = content.text_frame
    tf.text = "• Sistema de agentes especializados\n• Arquitectura escalable\n• Problemas identificados en auditoría"
    
    # Slide 3: Mejoras implementadas
    slide = prs.slides.add_slide(blank_layout)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title.text = "Mejoras Implementadas"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(5))
    tf = content.text_frame
    tf.text = "• FIX: matching logic en orquestador.js\n• Security: execFile() en server.js\n• Security: auth en endpoints Django"
    
    prs.save('/home/cesar/Documentos/agente-multiagente/presentacion_final.pptx')
    return '/home/cesar/Documentos/agente-multiagente/presentacion_final.pptx'

if __name__ == '__main__':
    ruta = crear_pptx()
    print(f"Creado: {ruta}")